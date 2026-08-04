# Document Agent
# 支持生成 Word / Excel / PPT 文件

import uuid
from pathlib import Path

from docx import Document as WordDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = PROJECT_ROOT / "storage/outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def generate_word(title: str, content: str) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    doc = WordDocument()
    doc.add_heading(title, level=1)
    for line in content.splitlines():
        doc.add_paragraph(line)
    path = OUTPUT_DIR / f"{uuid.uuid4()}.docx"
    doc.save(path)
    return path


def generate_excel(title: str, content: str) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = title[:31] or "Sheet"
    for row_index, line in enumerate(content.splitlines(), start=1):
        for col_index, cell in enumerate(line.split("|"), start=1):
            ws.cell(row=row_index, column=col_index, value=cell.strip())
    path = OUTPUT_DIR / f"{uuid.uuid4()}.xlsx"
    wb.save(path)
    return path


def generate_ppt(title: str, outline: str) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = outline
    for line in outline.splitlines():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = line[:80] or "内容"
    path = OUTPUT_DIR / f"{uuid.uuid4()}.pptx"
    prs.save(path)
    return path


def run_document_agent(
    question: str,
    doc_type: str | None = None
):
    doc_type = (doc_type or "word").lower()
    title = question[:30] or "生成文档"
    content = f"{question}\n\n由 Personal Office Agent 自动生成。"

    if doc_type == "ppt":
        path = generate_ppt(title, content)
    elif doc_type == "excel":
        path = generate_excel(title, content)
    else:
        path = generate_word(title, content)

    return {
        "answer": f"已生成文档：{path}",
        "sources": [path.name],
        "tool_result": str(path)
    }
