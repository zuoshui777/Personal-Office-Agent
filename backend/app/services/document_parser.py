# 读取用户上传的文件，把里面的文字提取出来
# 支持 PDF、DOCX、TXT、MD 文件解析


import os
import subprocess
import sys
from pathlib import Path

import fitz
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


SUPPORTED_SUFFIXES = {
    ".pdf",
    ".docx",
    ".txt",
    ".md",
    ".xlsx",
    ".pptx",
    ".py",
    ".java",
    ".csv",
    ".json",
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
    ".bmp"
}



# =========================
# PDF解析
# =========================

def parse_pdf(file_path: str) -> str:

    if not os.path.exists(file_path):
        raise FileNotFoundError(
            "PDF文件不存在"
        )


    text = ""


    # 自动关闭文件
    with fitz.open(
        os.path.abspath(file_path)
    ) as pdf:


        for page in pdf:

            page_text = page.get_text()

            if page_text.strip():

                text += page_text + "\n"


    if not text.strip():

        raise ValueError(
            "PDF没有可提取文本"
        )


    return text



# =========================
# DOCX解析
# =========================

def parse_docx(file_path: str) -> str:


    doc = Document(
        os.path.abspath(file_path)
    )


    text = []


    # 普通段落

    for paragraph in doc.paragraphs:

        if paragraph.text.strip():

            text.append(
                paragraph.text
            )


    # 表格内容

    for table in doc.tables:

        for row in table.rows:

            for cell in row.cells:

                if cell.text.strip():

                    text.append(
                        cell.text
                    )


    result = "\n".join(text)


    if not result.strip():

        raise ValueError(
            "DOCX没有文本内容"
        )


    return result



# =========================
# TXT解析
# =========================

def parse_txt(file_path: str) -> str:


    encodings = [
        "utf-8",
        "gbk",
        "gb2312"
    ]


    for encoding in encodings:

        try:

            with open(
                file_path,
                "r",
                encoding=encoding
            ) as f:

                return f.read()


        except UnicodeDecodeError:

            continue



    raise ValueError(
        "TXT编码无法识别"
    )



# =========================
# Markdown解析
# =========================

def parse_md(file_path: str) -> str:


    with open(
        file_path,
        "r",
        encoding="utf-8"
    ) as f:

        return f.read()


def parse_xlsx(file_path: str) -> str:
    wb = load_workbook(file_path, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            lines.append(" | ".join(
                "" if cell is None else str(cell)
                for cell in row
            ))
    return "\n".join(lines)


def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    lines = []
    for index, slide in enumerate(prs.slides, start=1):
        lines.append(f"Slide {index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    return "\n".join(lines)


def parse_code_file(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def parse_image(file_path: str) -> str:
    script = Path.home() / "plugins/mimo-vision/scripts/describe_image.py"
    if not script.exists():
        return "图片已上传，需要 MIMO Vision 插件识别。"
    result = subprocess.run(
        [
            sys.executable,
            str(script),
            str(file_path),
            "--prompt",
            "请完整识别这张图片中的文字和内容。"
        ],
        capture_output=True,
        text=True,
        timeout=180
    )
    if result.returncode == 0:
        return result.stdout
    return f"图片识别失败：{result.stderr.strip()}"



# =========================
# 自动选择解析器
# =========================

def parse_document(
    file_path: str
) -> str:


    suffix = os.path.splitext(
        file_path
    )[1].lower()


    parser_map = {

        ".pdf": parse_pdf,

        ".docx": parse_docx,

        ".txt": parse_txt,

        ".md": parse_md,

        ".xlsx": parse_xlsx,

        ".pptx": parse_pptx,

        ".py": parse_code_file,

        ".java": parse_code_file,

        ".csv": parse_code_file,

        ".json": parse_code_file,

        ".png": parse_image,

        ".jpg": parse_image,

        ".jpeg": parse_image,

        ".gif": parse_image,

        ".webp": parse_image,

        ".bmp": parse_image

    }


    if suffix not in parser_map:

        raise ValueError(
            f"暂不支持 {suffix} 文件"
        )


    return parser_map[suffix](
        file_path
    )
